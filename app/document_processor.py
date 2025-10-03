from typing import List, Optional, Union
from io import BytesIO
import pandas as pd
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

class DocumentProcessor:
    """
    A class to handle text extraction from various document formats including
    PDF, Word, PowerPoint, Excel, and text files.
    """
    
    @staticmethod
    def extract_from_pdf(file: BytesIO) -> str:
        """
        Extracts text from a PDF file.
        
        Args:
            file (BytesIO): PDF file in binary format
            
        Returns:
            str: Extracted text from the PDF
        """
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text

    @staticmethod
    def extract_from_docx(file: BytesIO) -> str:
        """
        Extracts text from a Word document.
        
        Args:
            file (BytesIO): Word document in binary format
            
        Returns:
            str: Extracted text from the document
        """
        doc = Document(file)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return "\n".join(text)

    @staticmethod
    def extract_from_pptx(file: BytesIO) -> str:
        """
        Extracts text from a PowerPoint presentation.
        
        Args:
            file (BytesIO): PowerPoint file in binary format
            
        Returns:
            str: Extracted text from the presentation
        """
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def extract_from_excel(file: BytesIO, sheet_name: Optional[Union[str, int]] = 0) -> str:
        """
        Extracts text from an Excel file (xlsx or csv).
        
        Args:
            file (BytesIO): Excel file in binary format
            sheet_name: Sheet name or index (default: 0)
            
        Returns:
            str: Extracted text from the spreadsheet
        """
        try:
            # Try reading as xlsx
            df = pd.read_excel(file, sheet_name=sheet_name)
        except:
            # If fails, try reading as csv
            file.seek(0)  # Reset file pointer
            df = pd.read_csv(file)
        
        # Convert DataFrame to string representation
        return df.to_string()

    @staticmethod
    def extract_from_txt(file: BytesIO) -> str:
        """
        Extracts text from a plain text file.
        
        Args:
            file (BytesIO): Text file in binary format
            
        Returns:
            str: Extracted text from the file
        """
        return file.read().decode('utf-8')

    def extract_text(self, file: BytesIO, file_type: str) -> str:
        """
        Main method to extract text from any supported document type.
        
        Args:
            file (BytesIO): File in binary format
            file_type (str): Type of the file (pdf, docx, pptx, xlsx, csv, txt)
            
        Returns:
            str: Extracted text from the document
        """
        extractors = {
            'pdf': self.extract_from_pdf,
            'docx': self.extract_from_docx,
            'pptx': self.extract_from_pptx,
            'xlsx': self.extract_from_excel,
            'csv': self.extract_from_excel,
            'txt': self.extract_from_txt
        }
        
        if file_type.lower() not in extractors:
            raise ValueError(f"Unsupported file type: {file_type}")
            
        return extractors[file_type.lower()](file)